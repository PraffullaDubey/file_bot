from flask import Flask, render_template, request, jsonify
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
import subprocess
import sys

app = Flask(__name__)

SUPPORTED_EXT = [".pdf", ".docx", ".pptx", ".txt", ".xlsx"]
file_index = {}
folder_indexed = False  # Track if folder has been indexed
folder_path = ""  # Store the folder path

def extract_text(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            reader = PdfReader(filepath)
            text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif ext == ".docx":
            doc = Document(filepath)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == ".pptx":
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == ".txt":
            with open(filepath, "r", encoding="utf-8") as f:
                text = f.read()
        elif ext == ".xlsx":
            wb = openpyxl.load_workbook(filepath, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"
    except Exception as e:
        print(f"Failed to read {filepath}: {e}")
    return text.lower()

def index_files(root_dir):
    index = {}
    for dirpath, _, filenames in os.walk(root_dir):
        for file in filenames:
            filepath = os.path.join(dirpath, file)
            if os.path.splitext(file)[1].lower() in SUPPORTED_EXT:
                print(f"Indexing: {filepath}")
                index[filepath] = extract_text(filepath)
    return index

def query_files(query, index):
    query = query.lower()
    results = []
    for path, content in index.items():
        if query in content:
            results.append(path)
    return results

@app.route("/")
def index():
    return render_template("chatbot.html")

@app.route("/index_folder", methods=["POST"])
def index_folder():
    global folder_indexed, folder_path, file_index
    folder = request.json.get("folder_path").strip()
    if os.path.isdir(folder):
        folder_path = folder  # Store folder path
        print("Indexing files...")
        file_index = index_files(folder)
        folder_indexed = True
        print(f"Indexed {len(file_index)} files.")
        return jsonify({"message": "Folder indexed successfully!"})
    return jsonify({"message": "Invalid folder path!"})

@app.route("/query_files", methods=["POST"])
def query_files_route():
    query = request.json.get("query")
    if not folder_indexed:
        return jsonify({"files": [], "message": "Folder has not been indexed. Please index a folder first."})
    
    matches = query_files(query, file_index)
    return jsonify({"files": matches})

@app.route("/open_file", methods=["POST"])
def open_file_route():
    file_path = request.json.get("file_path")
    try:
        if os.path.isfile(file_path):
            open_file(file_path)  # Use OS to open the file
            return jsonify({"message": f"Opening file: {file_path}"})
        else:
            return jsonify({"message": "File not found."})
    except Exception as e:
        return jsonify({"message": f"Error opening file: {e}"})

def open_file(filepath):
    try:
        if os.name == 'nt':  # Windows
            os.startfile(filepath)
        elif os.name == 'posix':  # macOS/Linux
            subprocess.call(['open' if sys.platform == 'darwin' else 'xdg-open', filepath])
    except Exception as e:
        print(f"Error opening file: {e}")

if __name__ == "__main__":
    app.run(debug=True)
