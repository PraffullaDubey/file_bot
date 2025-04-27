import os
from fpdf import FPDF
from docx import Document
from pptx import Presentation

output_dir = "sample_files"
os.makedirs(output_dir, exist_ok=True)

pdf = FPDF()
pdf.add_page()
pdf.set_font("Arial", size=12)
pdf.multi_cell(0, 10, "This is a sample PDF file.\nIt contains some sample text for testing purposes.\nSummary of automation plan.")
pdf_file_path = os.path.join(output_dir, "project_overview.pdf")
pdf.output(pdf_file_path)
print(f"PDF file created: {pdf_file_path}")

doc = Document()
doc.add_heading("Sample Word Document", level=1)
doc.add_paragraph("This is a sample Word document.\nIt contains some sample text for testing purposes.")
doc.add_paragraph("Next steps on Machine learning.")
doc_file_path = os.path.join(output_dir, "meeting_notes.docx")
doc.save(doc_file_path)
print(f"Word document created: {doc_file_path}")

ppt = Presentation()
slide_layout = ppt.slide_layouts[1]
slide = ppt.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Sample PowerPoint Presentation"
content = slide.shapes.placeholders[1]
content.text = "This is a sample PowerPoint presentation.\nIt contains some sample text for testing purposes regarding project plan."
ppt_file_path = os.path.join(output_dir, "project_plan.pptx")
ppt.save(ppt_file_path)
print(f"PowerPoint presentation created: {ppt_file_path}")